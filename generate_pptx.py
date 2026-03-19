import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)
CONTENT_W = SLIDE_W - 2 * MARGIN
HEADER_H = Inches(0.75)
FOOTER_H = Inches(0.4)

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
YELLOW = RGBColor(255, 255, 0)
SLATE = RGBColor(102, 102, 102)
BORDER = RGBColor(220, 220, 220)
LIGHT_BG = RGBColor(238, 238, 238)
BADGE_COL = RGBColor(42, 208, 169)

STATUS_COLORS = {
    "available":     (RGBColor(106, 151, 40),  RGBColor(231, 242, 217)),
    "partial":       (RGBColor(217, 115, 11),  RGBColor(254, 236, 209)),
    "in-progress":   (RGBColor(6, 150, 215),   RGBColor(205, 234, 247)),
    "coming-soon":   (RGBColor(6, 150, 215),   RGBColor(205, 234, 247)),
    "not-available": (RGBColor(221, 34, 34),   RGBColor(250, 237, 237)),
    "not-planned":   (RGBColor(153, 153, 153), RGBColor(238, 238, 238)),
}
STATUS_LABELS = {
    "available": "Available", "partial": "Partial", "in-progress": "In Progress",
    "coming-soon": "Coming Soon", "not-available": "Not Available", "not-planned": "Not Planned",
}

COL_FEAT = 0.35
COL_STAT = 0.15
COL_NOTE = 0.50


def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def tbox(slide, l, t, w, h, text, sz=10, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Arial"
    p.alignment = align
    return tb


def add_header(slide, title, subtitle):
    rect(slide, 0, 0, SLIDE_W, HEADER_H, BLACK)
    tbox(slide, MARGIN, Inches(0.12), CONTENT_W, Inches(0.35), title, sz=20, color=WHITE)
    tbox(slide, MARGIN, Inches(0.45), CONTENT_W, Inches(0.25), subtitle, sz=10, color=YELLOW)


def add_footer(slide):
    rect(slide, 0, SLIDE_H - FOOTER_H, SLIDE_W, FOOTER_H, BLACK)
    tbox(slide, MARGIN, SLIDE_H - FOOTER_H + Inches(0.08), CONTENT_W, Inches(0.25),
         "Last updated: March 2026", sz=8, color=RGBColor(180, 180, 180))
    tbox(slide, MARGIN, SLIDE_H - FOOTER_H + Inches(0.08), CONTENT_W, Inches(0.25),
         "\u00a9 2026 Autodesk. All rights reserved.  www.autodesk.com",
         sz=8, color=RGBColor(180, 180, 180), align=PP_ALIGN.RIGHT)


def draw_status_pill(slide, left, top, status):
    fg, bg = STATUS_COLORS.get(status, STATUS_COLORS["not-available"])
    label = STATUS_LABELS.get(status, "Unknown")
    pw = Inches(len(label) * 0.055 + 0.3)
    pill = rect(slide, left, top, pw, Inches(0.18), bg)
    dot = slide.shapes.add_shape(9, left + Inches(0.07), top + Inches(0.06), Inches(0.06), Inches(0.06))
    dot.fill.solid()
    dot.fill.fore_color.rgb = fg
    dot.line.fill.background()
    pill.text_frame.paragraphs[0].text = "       " + label
    pill.text_frame.paragraphs[0].font.size = Pt(7)
    pill.text_frame.paragraphs[0].font.color.rgb = fg
    pill.text_frame.paragraphs[0].font.bold = True
    pill.text_frame.paragraphs[0].font.name = "Arial"
    pill.text_frame.word_wrap = False


def draw_badge(slide, left, top, text):
    bw = Inches(len(text) * 0.055 + 0.15)
    b = rect(slide, left, top, bw, Inches(0.14), BADGE_COL)
    b.text_frame.paragraphs[0].text = text.upper()
    b.text_frame.paragraphs[0].font.size = Pt(5.5)
    b.text_frame.paragraphs[0].font.color.rgb = WHITE
    b.text_frame.paragraphs[0].font.bold = True
    b.text_frame.paragraphs[0].font.name = "Arial"
    b.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    b.text_frame.word_wrap = False


def draw_table(slide, table, y):
    fw = int(CONTENT_W * COL_FEAT)
    sw = int(CONTENT_W * COL_STAT)
    nw = int(CONTENT_W * COL_NOTE)
    rh = Inches(0.30)

    tbox(slide, MARGIN, y, CONTENT_W, Inches(0.22), table["title"], sz=11, color=BLACK, bold=True)
    y += Inches(0.24)

    desc = table.get("description", "")
    if desc:
        dl = max(1, len(desc) // 130 + 1)
        dh = Inches(0.14 * dl + 0.04)
        tbox(slide, MARGIN, y, CONTENT_W, dh, desc, sz=7, color=SLATE)
        y += dh + Inches(0.04)

    rect(slide, MARGIN, y, CONTENT_W, Inches(0.22), LIGHT_BG)
    rect(slide, MARGIN, y + Inches(0.22), CONTENT_W, Emu(9525), BORDER)
    for label, cl in [("FEATURE", MARGIN), ("STATUS", MARGIN + fw), ("NOTES", MARGIN + fw + sw)]:
        tbox(slide, cl + Inches(0.08), y + Inches(0.03), Inches(1.5), Inches(0.18),
             label, sz=7, color=SLATE, bold=True)
    y += Inches(0.22)

    for feat in table["features"]:
        notes = feat.get("notes", "")
        nl = max(1, len(notes) // 70 + 1) if notes else 1
        crh = max(rh, Inches(0.15 * nl + 0.10))

        tbox(slide, MARGIN + Inches(0.08), y + Inches(0.04), fw - Inches(0.15), crh,
             feat["name"], sz=8, color=BLACK)

        badge_text = feat.get("badge") or ("New" if feat.get("isNew") else "")
        if badge_text:
            name_approx_w = Inches(len(feat["name"]) * 0.048 + 0.12)
            if name_approx_w < fw - Inches(0.4):
                draw_badge(slide, MARGIN + Inches(0.08) + name_approx_w, y + Inches(0.06), badge_text)

        draw_status_pill(slide, MARGIN + fw + Inches(0.08), y + Inches(0.05), feat.get("status", "not-available"))

        if notes:
            tbox(slide, MARGIN + fw + sw + Inches(0.08), y + Inches(0.04),
                 nw - Inches(0.15), crh, notes, sz=7, color=SLATE)

        rect(slide, MARGIN, y + crh, CONTENT_W, Emu(6350), BORDER)
        y += crh

    return y


def main():
    with open("data.json", "r") as f:
        data = json.load(f)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    usable_bottom = SLIDE_H - FOOTER_H - Inches(0.15)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, data["title"], data["subtitle"])
    add_footer(slide)
    y = HEADER_H

    intro = data.get("introduction", "")
    if intro:
        lines = intro.split("\n")
        combined = "\n".join(lines)
        lc = len(lines) + sum(len(l) // 140 for l in lines)
        ih = Inches(0.16 * lc + 0.12)
        rect(slide, 0, y, SLIDE_W, ih, RGBColor(248, 248, 248))
        rect(slide, 0, y + ih, SLIDE_W, Emu(9525), BORDER)
        tbox(slide, MARGIN, y + Inches(0.06), CONTENT_W, ih - Inches(0.06), combined, sz=7.5, color=SLATE)
        y += ih + Inches(0.1)

    for i, table in enumerate(data["tables"]):
        est_rows = len(table["features"])
        est_h = Inches(0.5 + est_rows * 0.32)

        if y + est_h > usable_bottom and y > HEADER_H + Inches(1):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_header(slide, data["title"], data["subtitle"])
            add_footer(slide)
            y = HEADER_H + Inches(0.15)

        y = draw_table(slide, table, y)
        y += Inches(0.18)

    out = "BC_Pro_to_Forma_Bid.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
